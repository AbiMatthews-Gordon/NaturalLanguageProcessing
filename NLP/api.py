import nltk
import PyPDF2
from flask import Flask, jsonify, request, make_response
from werkzeug import utils as wz_utils
import uuid

import os
import re

from PyPDF2 import PdfReader
from pptx import Presentation
import docx2txt
import pyttsx3

from models.lexical_analyzer import LexicalAnalyser
from models.parser import Parser

app = Flask(__name__)

# download nltk requirements
def nltk_downloads():
    nltk.download('wordnet')
    nltk.download('stopwords')
    nltk.download('maxent_ne_chunker')
    nltk.download('words')


# *********** API ENDPOINTS **************
# ******************************************
# post request for the convert_text method


@app.get("/")
def index():
    return jsonify({"Message": "Welcome to my Speech Engine"})


# post request for the enter_text method
@app.post("/text")
def enter_text():

    # Get text sent by client
    data = request.get_json()
    text = data["text"]
    # print(data)

    # call Convert(text)
    result = convert_text(text)

    return result


@app.post('/upload_file')
def upload_file():

    file = request.files['file']
    file.save(wz_utils.secure_filename(file.filename))

    # handle any exception that may occur while the file is being processed
    try:
        result = handle_file(wz_utils.secure_filename(file.filename))
        return jsonify(result)
    except Exception as ex:
        resp = make_response("An error occurred while processing file", 400) # client gets this msg
        return resp


# ************** UTILITIES ***************
# ****************************************

# file handler process file based on file type
def handle_file(filename):

    # getting the location of folder
    dirpath = os.path.dirname(os.path.realpath(__file__))
        # print
        # check which file type - pdf, docx, pptx, txt
    if filename.endswith('.pdf'):
        text = "";
        reader = PdfReader(dirpath + "\\" + filename)
        number_of_pages = len(reader.pages)
        for page in reader.pages:
            print(page.extract_text())
           # page = reader.pages[0]
            text = text + page.extract_text()

        result = convert_text(text)
        # print(result)
        return result

    elif filename.endswith('.docx'):
        text = docx2txt.process(dirpath + "\\" + filename)
        result = convert_text(text)
        return result

    elif filename.endswith('.pptx'):
        with open(dirpath + "\\" + filename, "rb") as file:
            pptFile = Presentation(file)
            text = ""
            for slides in pptFile.slides:
                for shape in slides.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = text + run.text

            result = convert_text(text)
            # print(result)
            return result

    elif filename.endswith('.txt'):
        with open(dirpath + "\\" + filename) as file:

            text = file.read()
            result = convert_text(text)
            # print(result)
            return result
    else:
        raise Exception("Invalid file")


def convert_text(text):
    result = None
    parse_tree = None
    alpanu = re.compile('\w+').findall(text)  # regex to check if file has letters and/or numbers

    # check if string is valid - must have length >0, numbers &/or letters
    if len(text) > 0 and alpanu:
        # call lexical analyzer
        tokens = LexicalAnalyser.perform_lexical_analysis(text)
        # create folder to save parsar and audio results
        foldername = str(uuid.uuid4())
        os.mkdir("static\\" + foldername)
        # call parser on results

        parse_tree = Parser.generate_parser_tree(tokens.get("pos_sentences"), foldername)

        print("-------------------------")
        print("Processed text:/n ")
        print(tokens.get("process_text"))
        print("-------------------------")

        # generate audio
        generate_audio(tokens.get("process_text"), foldername)
        parse_tree.update({"tokens": tokens.get("pos_sentences"), "id": foldername, "audio_file": ("/static/" + foldername + "/audio.mp3")})
        return parse_tree


def generate_audio(text, foldername):

    # create the object of the library
    engine = pyttsx3.init()

    # tell the object what to save to the file
    engine.save_to_file(text, "static\\" + foldername +'\\audio.mp3')

    # set the object to runAndWait
    engine.runAndWait()


# output = convert_text("The cow jumped over the moon.")
# print(output)
nltk_downloads()
app.run(debug=True)
